# Document processing service for text extraction from various file formats
import os
import io
import logging
import mimetypes
from typing import Dict, Any, Optional, Tuple
from pathlib import Path
import tempfile
import asyncio

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Service for extracting text content from various document formats"""
    
    def __init__(self):
        self.supported_types = {
            'text/plain': self._extract_text,
            'text/markdown': self._extract_text,
            'text/csv': self._extract_text,
            'application/json': self._extract_text,
            'text/html': self._extract_html,
            'application/pdf': self._extract_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._extract_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx,
            'application/msword': self._extract_doc,
            'application/vnd.ms-excel': self._extract_xls,
            'application/vnd.ms-powerpoint': self._extract_ppt,
        }
    
    def is_supported(self, content_type: str) -> bool:
        """Check if the content type is supported for text extraction"""
        return content_type in self.supported_types
    
    async def extract_text(self, file_content: bytes, content_type: str, filename: str = "") -> Tuple[str, Dict[str, Any]]:
        """
        Extract text content from file bytes
        
        Args:
            file_content: Raw file bytes
            content_type: MIME type of the file
            filename: Original filename (optional, used for type detection)
        
        Returns:
            Tuple of (extracted_text, metadata)
        """
        try:
            logger.info(f"Extracting text from {content_type} file: {filename}")
            
            # Auto-detect content type if needed
            if not content_type or content_type == 'application/octet-stream':
                detected_type = self._detect_content_type(file_content, filename)
                if detected_type:
                    content_type = detected_type
                    logger.info(f"Auto-detected content type: {content_type}")
            
            if not self.is_supported(content_type):
                logger.warning(f"Unsupported content type: {content_type}")
                return "", {"error": f"Unsupported file type: {content_type}"}
            
            # Extract text using appropriate method
            extractor = self.supported_types[content_type]
            text_content = await extractor(file_content)
            
            # Generate metadata
            metadata = {
                "content_type": content_type,
                "filename": filename,
                "character_count": len(text_content),
                "word_count": len(text_content.split()) if text_content else 0,
                "extraction_success": True
            }
            
            logger.info(f"Successfully extracted {len(text_content)} characters from {filename}")
            return text_content, metadata
            
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {str(e)}")
            return "", {"error": str(e), "extraction_success": False}
    
    def _detect_content_type(self, file_content: bytes, filename: str) -> Optional[str]:
        """Auto-detect content type from file extension or content"""
        if filename:
            # Try to detect from extension
            content_type, _ = mimetypes.guess_type(filename)
            if content_type:
                return content_type
        
        # Try to detect from content (basic detection)
        if file_content.startswith(b'%PDF'):
            return 'application/pdf'
        elif file_content.startswith(b'PK'):  # ZIP-based formats (docx, xlsx, etc.)
            if b'word/' in file_content[:1000]:
                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif b'xl/' in file_content[:1000]:
                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif b'ppt/' in file_content[:1000]:
                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        return None
    
    async def _extract_text(self, file_content: bytes) -> str:
        """Extract text from plain text files"""
        try:
            # Try UTF-8 first, fallback to latin-1
            try:
                return file_content.decode('utf-8')
            except UnicodeDecodeError:
                return file_content.decode('latin-1', errors='ignore')
        except Exception as e:
            logger.error(f"Error extracting plain text: {str(e)}")
            return ""
    
    async def _extract_html(self, file_content: bytes) -> str:
        """Extract text from HTML files"""
        try:
            # Try to import BeautifulSoup for better HTML parsing
            try:
                from bs4 import BeautifulSoup
                html_content = file_content.decode('utf-8', errors='ignore')
                soup = BeautifulSoup(html_content, 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator=' ', strip=True)
            except ImportError:
                # Fallback to simple HTML tag removal
                import re
                html_content = file_content.decode('utf-8', errors='ignore')
                clean_text = re.sub('<[^<]+?>', '', html_content)
                return clean_text.strip()
        except Exception as e:
            logger.error(f"Error extracting HTML text: {str(e)}")
            return ""
    
    async def _extract_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF files"""
        try:
            # Try multiple PDF extraction methods
            text_content = ""
            
            # Method 1: PyPDF2
            try:
                import PyPDF2
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                for page in pdf_reader.pages:
                    text_content += page.extract_text() + "\n"
                
                if text_content.strip():
                    return text_content
            except ImportError:
                logger.warning("PyPDF2 not available for PDF extraction")
            except Exception as e:
                logger.warning(f"PyPDF2 extraction failed: {str(e)}")
            
            # Method 2: pdfplumber (if available)
            try:
                import pdfplumber
                with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_content += page_text + "\n"
                
                if text_content.strip():
                    return text_content
            except ImportError:
                logger.warning("pdfplumber not available for PDF extraction")
            except Exception as e:
                logger.warning(f"pdfplumber extraction failed: {str(e)}")
            
            # If no text extracted, return placeholder
            if not text_content.strip():
                return "[PDF content - text extraction not available or document contains only images]"
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            return "[PDF content - extraction failed]"
    
    async def _extract_docx(self, file_content: bytes) -> str:
        """Extract text from Word documents (.docx)"""
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(file_content))
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.warning("python-docx not available for DOCX extraction")
            return "[Word document - python-docx library required for text extraction]"
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            return "[Word document - extraction failed]"
    
    async def _extract_xlsx(self, file_content: bytes) -> str:
        """Extract text from Excel spreadsheets (.xlsx)"""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None and str(cell_value).strip():
                            row_text.append(str(cell_value))
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.warning("openpyxl not available for XLSX extraction")
            return "[Excel spreadsheet - openpyxl library required for text extraction]"
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {str(e)}")
            return "[Excel spreadsheet - extraction failed]"
    
    async def _extract_pptx(self, file_content: bytes) -> str:
        """Extract text from PowerPoint presentations (.pptx)"""
        try:
            from pptx import Presentation
            
            presentation = Presentation(io.BytesIO(file_content))
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_content.append(f"=== Slide {slide_num} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.warning("python-pptx not available for PPTX extraction")
            return "[PowerPoint presentation - python-pptx library required for text extraction]"
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            return "[PowerPoint presentation - extraction failed]"
    
    async def _extract_doc(self, file_content: bytes) -> str:
        """Extract text from legacy Word documents (.doc)"""
        try:
            # This requires antiword or similar tool, which is complex to install
            # For now, return a placeholder
            return "[Legacy Word document (.doc) - conversion to .docx recommended for text extraction]"
        except Exception as e:
            logger.error(f"Error extracting DOC text: {str(e)}")
            return "[Legacy Word document - extraction failed]"
    
    async def _extract_xls(self, file_content: bytes) -> str:
        """Extract text from legacy Excel files (.xls)"""
        try:
            import xlrd
            
            workbook = xlrd.open_workbook(file_contents=file_content)
            text_content = []
            
            for sheet_name in workbook.sheet_names():
                sheet = workbook.sheet_by_name(sheet_name)
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        if cell_value and str(cell_value).strip():
                            row_values.append(str(cell_value))
                    if row_values:
                        text_content.append(" | ".join(row_values))
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.warning("xlrd not available for XLS extraction")
            return "[Legacy Excel file (.xls) - xlrd library required or convert to .xlsx]"
        except Exception as e:
            logger.error(f"Error extracting XLS text: {str(e)}")
            return "[Legacy Excel file - extraction failed]"
    
    async def _extract_ppt(self, file_content: bytes) -> str:
        """Extract text from legacy PowerPoint files (.ppt)"""
        try:
            # Legacy PPT extraction is complex, recommend conversion
            return "[Legacy PowerPoint file (.ppt) - conversion to .pptx recommended for text extraction]"
        except Exception as e:
            logger.error(f"Error extracting PPT text: {str(e)}")
            return "[Legacy PowerPoint file - extraction failed]"

# Global document processor instance
document_processor = DocumentProcessor()